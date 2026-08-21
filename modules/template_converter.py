"""
Multi-format certificate template converter.

Accepts SVG, PDF, DOCX, PPTX, PNG, JPG/JPEG template uploads and produces
a single, consistent output: an editable SVG certificate template with
{{participant_name}}, {{course_name}}, {{description}}, {{issue_date}},
{{certificate_id}} text placeholders and a QR code slot - the same format
every certificate in this app is generated from.

Strategy per format:
  - SVG (with real {{...}} tokens) : used directly (namespaces + QR slot ensured).
  - SVG (no real tokens found)     : rasterized and wrapped like a photo.
  - PNG/JPG      : wrapped as an SVG background.
  - PDF          : first page rasterized, then wrapped.
  - DOCX / PPTX  : the largest embedded picture is extracted and wrapped.
"""

import os
import re
import tempfile

from modules.svg_certificate_gen import (
    wrap_raster_image_as_svg,
    ensure_svg_namespaces,
    ensure_qr_placeholder,
    PLACEHOLDER_KEYS,
)


class TemplateConversionError(Exception):
    """Raised when a template file can't be converted, with a user-facing message."""
    pass


def _svg_has_real_placeholders(svg_source):
    """
    True only when at least one known certificate placeholder token exists
    as real substitutable text (not just a random '{{' somewhere).
    """
    for key in PLACEHOLDER_KEYS:
        if "{{" + key + "}}" in svg_source:
            return True
    # Also accept common short aliases
    for alias in ("{{name}}", "{{date}}", "{{event}}"):
        if alias in svg_source:
            return True
    return False


def _flatten_svg_to_png(svg_path):
    """
    Rasterizes an SVG file to a PNG when it has no real editable text left.
    Uses WeasyPrint → PDF → PyMuPDF PNG.
    """
    try:
        from weasyprint import HTML
        import fitz  # PyMuPDF
    except ImportError as e:
        raise TemplateConversionError(
            f"Converting this SVG requires additional packages: {e}. "
            "Please run: pip install weasyprint pymupdf"
        )

    with open(svg_path, "r", encoding="utf-8", errors="replace") as f:
        svg_source = f.read()

    from modules.svg_certificate_gen import get_svg_dimensions
    width, height = get_svg_dimensions(svg_source)
    width, height = int(width) or 1200, int(height) or 800
    width = max(width, 100)
    height = max(height, 100)

    svg_source = ensure_svg_namespaces(svg_source)

    html_wrapper = f"""<!DOCTYPE html><html><head><meta charset="UTF-8"><style>
@page {{ size: {width}px {height}px; margin: 0; }}
html, body {{ margin: 0; padding: 0; width: {width}px; height: {height}px; }}
svg {{ width: {width}px; height: {height}px; display: block; }}
</style></head><body>{svg_source}</body></html>"""

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
        tmp_pdf_path = tmp_pdf.name

    tmp_png_path = None
    try:
        HTML(string=html_wrapper, base_url=os.path.dirname(svg_path) or ".").write_pdf(tmp_pdf_path)

        doc = fitz.open(tmp_pdf_path)
        page = doc[0]
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2))

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_png:
            tmp_png_path = tmp_png.name
        pixmap.save(tmp_png_path)
        doc.close()

    except Exception as e:
        raise TemplateConversionError(f"Could not render this SVG file: {e}")
    finally:
        if os.path.exists(tmp_pdf_path):
            try:
                os.remove(tmp_pdf_path)
            except OSError:
                pass

    return tmp_png_path


def convert_uploaded_template(file_path, extension):
    """
    Converts an uploaded template file of any supported type into SVG source text.

    Returns: (svg_source: str, warning: str or None)
    Raises: TemplateConversionError with a friendly message on failure.
    """
    extension = (extension or "").lower().lstrip(".")

    if extension == "svg":
        return _convert_svg(file_path)

    if extension in ("png", "jpg", "jpeg"):
        return wrap_raster_image_as_svg(file_path), None

    if extension == "pdf":
        return _convert_pdf(file_path)

    if extension == "docx":
        return _convert_docx(file_path)

    if extension == "pptx":
        return _convert_pptx(file_path)

    raise TemplateConversionError(
        f"Unsupported file type '.{extension}'. "
        "Please upload an SVG, PDF, DOCX, PPTX, PNG, or JPG file."
    )


def _convert_svg(file_path):
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        svg_source = f.read()

    if _svg_has_real_placeholders(svg_source):
        # Keep the author's layout; just make sure namespaces + QR slot exist.
        svg_source = ensure_svg_namespaces(svg_source)
        svg_source = ensure_qr_placeholder(svg_source)
        warning = None
        if 'data-placeholder="qr_code"' not in svg_source and "data-placeholder='qr_code'" not in svg_source:
            warning = (
                "Your SVG has name/date placeholders but no QR code slot. "
                "A QR placeholder was added in the bottom-right corner automatically."
            )
        return svg_source, warning

    # No real {{...}} tokens — text was likely flattened to outlines on export.
    png_path = _flatten_svg_to_png(file_path)

    try:
        wrapped_svg = wrap_raster_image_as_svg(png_path)
    finally:
        if png_path and os.path.exists(png_path):
            try:
                os.remove(png_path)
            except OSError:
                pass

    warning = (
        "This SVG file doesn't contain any editable placeholder text (its text "
        "appears to have been converted to outlines/shapes when it was exported, "
        "which is common from design tools like Illustrator, Figma, or Canva). "
        "It has been converted to an image-based design instead, with working "
        "placeholder fields (name, date, certificate ID, description, QR code) added "
        "on top. For a fully clean result, re-export the design with text kept as real, "
        "editable text (not outlines), using tokens like {{participant_name}}."
    )

    return wrapped_svg, warning


def _convert_pdf(file_path):
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise TemplateConversionError(
            "PDF template support requires the 'PyMuPDF' package. "
            "Please run: pip install pymupdf"
        )

    tmp_path = None
    try:
        doc = fitz.open(file_path)
        page = doc[0]
        has_real_text = bool(page.get_text().strip())

        matrix = fitz.Matrix(2, 2)
        pixmap = page.get_pixmap(matrix=matrix)

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        pixmap.save(tmp_path)
        doc.close()

    except Exception as e:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except OSError:
                pass
        raise TemplateConversionError(f"Could not read this PDF file: {e}")

    try:
        svg_source = wrap_raster_image_as_svg(tmp_path)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except OSError:
                pass

    if has_real_text:
        warning = (
            "This PDF's page was converted to an image-based certificate design. "
            "Editable placeholder fields (name, date, certificate ID, description, QR code) "
            "have been added on top."
        )
    else:
        warning = (
            "This PDF appears to be a scanned or rasterized page. "
            "Placeholder fields (name, date, certificate ID, description, QR code) "
            "have been added automatically so certificates can still be generated."
        )

    return svg_source, warning


def _extract_largest_image_from_docx(file_path):
    from docx import Document

    doc = Document(file_path)

    best_blob = None
    best_size = 0
    best_ext = "png"

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            blob = rel.target_part.blob
            if len(blob) > best_size:
                best_size = len(blob)
                best_blob = blob
                best_ext = "jpg" if "jpeg" in (rel.target_part.content_type or "") else "png"

    return best_blob, best_ext


def _convert_docx(file_path):
    try:
        from docx import Document  # noqa: F401
    except ImportError:
        raise TemplateConversionError(
            "DOCX template support requires the 'python-docx' package. "
            "Please run: pip install python-docx"
        )

    try:
        blob, ext = _extract_largest_image_from_docx(file_path)
    except Exception as e:
        raise TemplateConversionError(f"Could not read this Word document: {e}")

    if blob is None:
        raise TemplateConversionError(
            "No embedded image was found in this Word document. "
            "Please insert the certificate design as a picture, or upload PDF/PNG/JPG/SVG instead."
        )

    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
        tmp_path = tmp.name
        tmp.write(blob)

    try:
        svg_source = wrap_raster_image_as_svg(tmp_path)
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass

    warning = (
        "The background image was extracted from your Word document and converted to an "
        "editable certificate design. Other text/formatting in the document was not used."
    )

    return svg_source, warning


def _extract_largest_image_from_pptx(file_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)

    if not prs.slides:
        return None, "png"

    slide = prs.slides[0]

    best_blob = None
    best_area = 0
    best_ext = "png"

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            area = shape.width * shape.height
            if area > best_area:
                best_area = area
                best_blob = shape.image.blob
                best_ext = shape.image.ext or "png"

    return best_blob, best_ext


def _convert_pptx(file_path):
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        raise TemplateConversionError(
            "PPTX template support requires the 'python-pptx' package. "
            "Please run: pip install python-pptx"
        )

    try:
        blob, ext = _extract_largest_image_from_pptx(file_path)
    except Exception as e:
        raise TemplateConversionError(f"Could not read this PowerPoint file: {e}")

    if blob is None:
        raise TemplateConversionError(
            "No embedded image was found on the first slide of this presentation. "
            "Please insert the design as a picture, or upload PDF/PNG/JPG/SVG instead."
        )

    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
        tmp_path = tmp.name
        tmp.write(blob)

    try:
        svg_source = wrap_raster_image_as_svg(tmp_path)
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass

    warning = (
        "The background image was extracted from the first slide of your presentation and "
        "converted to an editable certificate design. Other slides/content were not used."
    )

    return svg_source, warning